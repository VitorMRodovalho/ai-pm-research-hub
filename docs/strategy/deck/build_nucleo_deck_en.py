#!/usr/bin/env python
# Núcleo IA & GP — executive pitch deck, ENGLISH edition (board PMI / chapter presidents / partners).
# Sibling of build_nucleo_deck.py (PT): SAME engine, translated copy. The Núcleo name is kept
# verbatim ("Núcleo IA & GP") as the dual-language brand; the cover glosses it in English.
# Brand: PMIGO via the official PMI Events template (clone + inject by shape name).
# Build = pptx -> PDF -> preview PNGs in ONE run. Guards: em-dash, bounds, boilerplate, divider-overlap.
import copy, subprocess, shutil
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

BASE = Path("/home/vitormrodovalho/projects/ai-pm-research-hub/docs/strategy/deck")
TEMPLATE = BASE / "build/pmi_template.pptx"
OUT = BASE / "Nucleo_IA_GP_Pitch_Executive_EN.pptx"
PREVIEW = BASE / "preview_en"
HUB = str(BASE / "assets/hub_spoke_en.png")
LOGO = "/home/vitormrodovalho/projects/19SGPL-PMIGO/MKT/Pending/pmi_chp_logo_goias_brazil_hrz_wht.png"

PURPLE=RGBColor(0x46,0x1D,0xA3); PURPLE_DK=RGBColor(0x36,0x17,0x7B); LILAC=RGBColor(0xA9,0x8E,0xEC)
TEAL=RGBColor(0x44,0x78,0x9B); BLUE=RGBColor(0x6C,0xBE,0xDE); RED=RGBColor(0xB8,0x37,0x13)
ORANGE=RGBColor(0xE0,0x61,0x1F); GREEN=RGBColor(0x3E,0x8E,0x5A); TAN=RGBColor(0xB2,0x94,0x78)
LIGHT=RGBColor(0xF6,0xF4,0xEF); DARK=RGBColor(0x24,0x20,0x16); WHITE=RGBColor(0xFF,0xFF,0xFF)
GRAY=RGBColor(0x60,0x60,0x60)
FONT="Aptos"
COVER, CONTENT = 10, 17

prs = Presentation(str(TEMPLATE))
SRC = list(prs.slides._sldIdLst)
slides = prs.slides
SW_IN = prs.slide_width / 914400; SH_IN = prs.slide_height / 914400
K = SW_IN / 13.333
CT = 2.75 / K
INJECTED = []

def IN(v): return Inches(v * K)
def FS(p): return Pt(round(p * K))

def _iter(shapes):
    for sh in shapes:
        if sh.shape_type == 6: yield from _iter(sh.shapes)
        else: yield sh

def clone(src_idx):
    src = slides[src_idx]
    dst = slides.add_slide(src.slide_layout)
    for ph in list(dst.shapes):
        ph._element.getparent().remove(ph._element)
    rid_map = {}
    for rId, rel in src.part.rels.items():
        if rel.reltype.endswith(("slideLayout", "notesSlide")): continue
        rid_map[rId] = (dst.part.rels.get_or_add_ext_rel(rel.reltype, rel.target_ref)
                        if rel.is_external else dst.part.relate_to(rel.target_part, rel.reltype))
    spTree = dst.shapes._spTree
    for child in src.shapes._spTree.iterchildren():
        if child.tag.split('}')[-1] in ("nvGrpSpPr", "grpSpPr"): continue
        el = copy.deepcopy(child)
        for node in el.iter():
            for a in list(node.attrib):
                if a.startswith('{http://schemas.openxmlformats.org/officeDocument/2006/relationships}'):
                    if node.get(a) in rid_map: node.set(a, rid_map[node.get(a)])
        spTree.append(el)
    return dst

def named(slide, name): return [s for s in _iter(slide.shapes) if s.name == name]

def delete(slide, *names):
    for n in names:
        for s in list(slide.shapes):
            if s.name == n:
                s._element.getparent().remove(s._element)

def fix_year(slide):
    for sh in _iter(slide.shapes):
        if sh.has_text_frame and "2025" in sh.text_frame.text:
            for p in sh.text_frame.paragraphs:
                for r in p.runs:
                    if "2025" in r.text: r.text = r.text.replace("2025", "2026")

def inject_keep(shape, text):
    tf = shape.text_frame
    runs = tf.paragraphs[0].runs
    size = bold = name = color = None
    if runs:
        f = runs[0].font; size, bold, name = f.size, f.bold, f.name
        try: color = f.color.rgb
        except Exception: color = None
    tf.clear()
    for i, ln in enumerate(text.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        r = p.add_run(); r.text = ln
        if size: r.font.size = size
        if bold is not None: r.font.bold = bold
        if name: r.font.name = name
        if color is not None:
            try: r.font.color.rgb = color
            except Exception: pass

def notes(slide, text): slide.notes_slide.notes_text_frame.text = text

def _track(slide, sp): INJECTED.append((slide, sp.left, sp.top, sp.width, sp.height))

def add_box(slide, left, top, w, h, header, lines, hcolor=PURPLE, fs=13, hfs=15, body=DARK):
    tb = slide.shapes.add_textbox(IN(left), IN(top), IN(w), IN(h))
    tf = tb.text_frame; tf.word_wrap = True
    first = True
    if header:
        p = tf.paragraphs[0]; r = p.add_run(); r.text = header
        r.font.bold = True; r.font.size = FS(hfs); r.font.name = FONT; r.font.color.rgb = hcolor
        p.space_after = Pt(8); first = False
    for ln in lines:
        p = tf.paragraphs[0] if first else tf.add_paragraph(); first = False
        r = p.add_run(); r.text = ln
        r.font.size = FS(fs); r.font.name = FONT; r.font.color.rgb = body
        p.space_after = Pt(6)
    _track(slide, tb)
    return tb

def add_table(slide, rows, top, left=0.62, width=12.1, fs=12, widths=None):
    nr, nc = len(rows), len(rows[0])
    gf = slide.shapes.add_table(nr, nc, IN(left), IN(top), IN(width), IN(0.4*nr))
    t = gf.table
    if widths:
        tot = sum(widths)
        for j, w in enumerate(widths): t.columns[j].width = IN(width*w/tot)
    for i, row in enumerate(rows):
        for j, val in enumerate(row):
            c = t.cell(i, j)
            for m in ("left","right"): setattr(c, f"margin_{m}", IN(0.08))
            for m in ("top","bottom"): setattr(c, f"margin_{m}", IN(0.03))
            c.fill.solid()
            if i == 0:
                c.fill.fore_color.rgb = PURPLE; fc, fb = WHITE, True
            else:
                c.fill.fore_color.rgb = WHITE if i % 2 else LIGHT; fc, fb = DARK, (j == 0)
            tf = c.text_frame; tf.word_wrap = True; tf.clear()
            p = tf.paragraphs[0]; r = p.add_run(); r.text = str(val)
            r.font.size = FS(fs); r.font.bold = fb; r.font.color.rgb = fc; r.font.name = FONT
    _track(slide, gf)
    return gf

def add_image(slide, path, top, width, left=None):
    if left is None: left = (13.333 - width) / 2
    pic = slide.shapes.add_picture(path, IN(left), IN(top), width=IN(width))
    _track(slide, pic)
    return pic

def _arrow(slide, left, top, color=LILAC, size=30):
    tb = slide.shapes.add_textbox(IN(left), IN(top), IN(0.9), IN(0.8))
    p = tb.text_frame.paragraphs[0]; r = p.add_run(); r.text = "→"
    r.font.size = FS(size); r.font.bold = True; r.font.name = FONT; r.font.color.rgb = color
    return tb

BOILER = ("TextBox 9", "Group 10")

def content(eyebrow, title, builder, note=""):
    s = clone(CONTENT)
    tb = named(s, "TextBox 14")
    if tb: inject_keep(tb[0], title)
    eb = named(s, "TextBox 13")
    if eb: inject_keep(eb[0], eyebrow)
    delete(s, *BOILER)
    fix_year(s)
    builder(s, CT)
    if note: notes(s, note)
    return s

# ===================== SPECS (English; no em-dash) =====================

# ---- 1 COVER (dual-language name: brand kept, English gloss in the subtitle) ----
cov = clone(COVER)
ct = named(cov, "TextBox 14"); cs = named(cov, "TextBox 15")
if ct: inject_keep(ct[0], "Núcleo IA & GP")
if cs: inject_keep(cs[0], "AI & PM Study and Research Hub  ·  The seam across PMI's "
                          "credential silos  ·  June 2026")
fix_year(cov)
try:
    cov.shapes.add_picture(LOGO, Inches(0.74), Inches(0.45), height=Inches(0.85))
except Exception as e: print("logo skip:", e)
attr = cov.shapes.add_textbox(Inches(0.74), Inches(10.55), Inches(15.0), Inches(0.45))
ap = attr.text_frame.paragraphs[0]; ar = ap.add_run()
ar.text = ("PMI, the PMI logo, PMI-CP, PMI-PMOCP, PMI-ACP, CSPP, CPMAI, PfMP, PgMP and PMI-PBA are "
           "marks of the Project Management Institute, Inc. Used under the chapter brand guidelines.")
ar.font.size = Pt(10); ar.font.name = FONT; ar.font.color.rgb = LILAC
notes(cov, "Opening. The Nucleo IA & GP (kept as the dual-language brand) is the seam across PMI's "
           "credential silos: a volunteer community that brings good people from every credential into "
           "research, development and networking, with AI as the thread. This deck shows the thesis (the "
           "horizontal PMI does not have by design), the three-axis model, the verticals and the ask per "
           "audience. Audience: PMI board, chapter presidents and vertical partners. PMIGO brand under "
           "chapter guidelines; marks attribution line in the footer.")

# ---- 2 THE PROBLEM ----
def b_prob(s, top):
    add_box(s, 0.62, top, 5.95, 4.3, "PMI lives in credential silos", [
        "Construction, PMO, Agile, Sustainability, Business.",
        "Each community in its own track; they rarely talk.",
        "Every partnership (GPM, Construction Ambassadors, PMOGA) becomes an isolated special case."],
        hcolor=PURPLE)
    add_box(s, 6.95, top, 5.75, 4.3, "AI is transversal to all of them", [
        "No credential escapes AI: it cuts across every silo.",
        "The ecosystem lacks a horizontal that sews the boundaries together.",
        "That is the Núcleo's role: volunteer to volunteer, a single purpose, good people across credentials."],
        hcolor=ORANGE)
content("The problem", "What PMI is missing is a horizontal", b_prob,
    "The problem is not a lack of communities, it is the lack of a seam between them. PMI is organized by "
    "credential silos; AI is transversal to all of them. Without an explicit horizontal, every "
    "partnership becomes a special case. The Nucleo is that horizontal by design: a volunteer community, "
    "one purpose, bringing good people together across credential boundaries, with AI as the thread.")

# ---- 3 THE TAILWIND -- killer slide ----
def b_tail(s, top):
    add_table(s, [
        ["Silo", "The PMI's own move (verified, Jun/2026)"],
        ["PMO", "PMI acquired PMO-CP from PMOGA (2023) and relaunched it as PMI-PMOCP, ISO-accredited (Feb/2026); PMOGA now lives at pmoga.pmi.org"],
        ["ESG / Green", "GPM-b evolves into CSPP, co-branded PMI + GPM, aligned to the P5 Standard (effective 5 Jun/2026)"],
        ["Agile", "Agile Alliance joined PMI (2026)"],
    ], top, widths=[1.4, 6.8], fs=13)
    add_box(s, 0.62, top+2.95, 12.1, 1.1,
        "The Núcleo is the human and research expression of that same integration. AI is the thread.", [],
        hcolor=PURPLE, hfs=15)
content("The tailwind", "We are not fighting the tide: we are naming it", b_tail,
    "This is the killer slide for the board. The integrating-hub thesis does not row against the "
    "institutional tide, it names the tide: PMI itself is absorbing credential communities into itself "
    "(PMO-CP from PMOGA in 2023; GPM-b becomes CSPP in Jun/2026; Agile Alliance joins in 2026). The "
    "Nucleo is the human and research expression of that integration, with AI as the thread. A practical "
    "reading of PMI:Next and M.O.R.E., not one more initiative.")

# ---- 4 THE MODEL -- hub-and-spoke ----
def b_model(s, top):
    add_image(s, HUB, top-0.15, 9.7)
    add_box(s, 0.62, top+4.35, 12.1, 1.0, "", [
        "Three orthogonal axes: Quadrant (what) x Tribe (who produces) x Vertical (who it lands with). "
        "The ladder shared by every vertical: PMIxAI Champion (open) → CPMAI Study Group → PMI-CPMAI.",
        "A vertical is a dock, never a silo: the tribe produces, the vertical distributes. No one owns the knowledge."],
        hcolor=PURPLE, fs=12, hfs=13, body=DARK)
content("The model", "AI at the center, the communities as spokes", b_model,
    "The visual is the pitch. Center: Nucleo + AI, the seam. Spokes: the verticals, each a PMI credential "
    "community. Ring: the Champion > CPMAI ladder, the spine shared by every spoke. The model has three "
    "orthogonal axes: quadrant (type of knowledge, the Nucleo's own taxonomy), tribe (who produces, Axis "
    "A) and vertical (who it distributes to, Axis B). The vertical does not duplicate the quadrant: it "
    "defines audience and packaging. Anti-silo principle: the tribe produces, the vertical distributes; "
    "no vertical owns knowledge, it is a dock.")

# ---- 5 THE LADDER ----
def b_ladder(s, top):
    y = top + 0.35; h = 3.0
    add_box(s, 0.62, y, 3.5, h, "1 · PMIxAI Champion", [
        "Recognition and badge, open and light.",
        "Entry point (Axis B).",
        "Already a platform primitive (award_champion)."], hcolor=BLUE)
    _arrow(s, 4.25, y+1.0)
    add_box(s, 4.85, y, 3.5, h, "2 · CPMAI Study Group", [
        "Preparatory, within a tribe (Axis A).",
        "The bridge from open to credential.",
        "Already modeled as an initiative (cpmai_*)."], hcolor=TEAL)
    _arrow(s, 8.5, y+1.0)
    add_box(s, 9.1, y, 3.6, h, "3 · PMI-CPMAI", [
        "The seam credential.",
        "Interfaces with ALL verticals.",
        "The common point between communities that would not otherwise talk."], hcolor=PURPLE)
content("The common ladder", "The credential seam is a ladder, not a single point", b_ladder,
    "The seam has a ladder that every vertical climbs. PMIxAI Champion (open, gamified, already exists as "
    "award_champion) prepares the ground; the CPMAI Study Group (Axis A, preparatory, already modeled) is "
    "the bridge; PMI-CPMAI is the seam credential that interfaces with all verticals. Central isomorphism: "
    "in the PMI ecosystem the silos are credentials and the seam is AI (CPMAI); inside the Nucleo the seam "
    "is the same mechanism, with the Champion > CPMAI ladder as the shared spine.")

# ---- 6 REACH & COVERAGE ----
def b_reach(s, top):
    add_box(s, 0.62, top, 5.95, 4.0, "Brazil + LatAm in the foreground", [
        "Active members (researchers, leaders, curators) already across several Brazilian states.",
        "15 PMI chapters in Brazil; intra-Brazil diversity is the priority.",
        "Framing LatAm is the geographic-expansion message, without spelling it out in copy."],
        hcolor=PURPLE)
    add_box(s, 6.95, top, 5.75, 4.0, "International presence = networking", [
        "The Núcleo already gathers people in Brazil, Portugal, Italy and the USA.",
        "Not a sparse world heatmap: a networking-access asset, named by country.",
        "Ambassadors are few (about 4); the active community is larger than that."],
        hcolor=TEAL)
    add_box(s, 0.62, top+4.15, 12.1, 0.55,
        "Governance-first (LGPD): aggregate by chapter, state or country. Zero PII; individual pins only with opt-in.",
        [], hcolor=ORANGE, hfs=12)
content("Reach & coverage", "We break the geographic silo, not only the credential one", b_reach,
    "Calibration as of Jun/2026: the Nucleo is not only ambassadors (about 4); active members are already "
    "across several Brazilian states plus Portugal, Italy and the USA. Intra-Brazil diversity is the "
    "priority, but the international presence is a networking asset and should appear, named by country, "
    "not as a sparse world blur. On a public page everything aggregates by chapter/state/country, zero "
    "PII, individual pins only with recorded opt-in. The PMAIrevolution map can be reused, re-projected to "
    "Brazil/LatAm.")

# ---- 7..11 VERTICALS ----
def vertical(eyebrow, title, dor, teses, timing, prova, hc, note, bordo=None):
    def b(s, top):
        bh = 3.4 if bordo else 4.1
        add_box(s, 0.62, top, 5.95, bh, "Audience & pain", dor, hcolor=hc)
        add_box(s, 6.95, top, 5.75, bh, "AI thesis", teses, hcolor=PURPLE)
        if bordo:
            add_box(s, 0.62, top+bh+0.15, 12.1, 1.55, "Already on board: leadership in place",
                    bordo, hcolor=ORANGE, hfs=14, fs=12.5)
        else:
            add_box(s, 0.62, top+4.2, 12.1, 0.6, "", [f"Timing: {timing}", f"Anchor proof: {prova}"],
                    fs=11.5, body=DARK)
    return content(eyebrow, title, b, note)

vertical("Vertical · Construction", "Giga-projects: PMI-CP and AI in megaprojects",
    ["Leaders of megaprojects and infrastructure (A/E/C).",
     "Long projects, massive documentation, schedule and cost risk.",
     "Endless contracts and RFIs; data fragmented across stakeholders."],
    ["Risk and schedule analysis in megaprojects.",
     "Reading and summarizing contracts, RFIs and submittals at scale.",
     "Digital twins and field data as decision input."],
    "the PMI motto 'Megaprojects demand mega skills' (megaproject = over US$1bn, multi-year); the "
    "Construction Ambassadors advocate safety, efficiency and sustainability (hub pmicp.us).",
    "AI study on a typical megaproject claim or risk + webinar with an ambassador.",
    ORANGE,
    "Construction vertical. Anchor credential PMI-CP; partner Global Construction Ambassadors. ALREADY "
    "HAS protagonists: Henrique Diniz (Brazil) and Fabricio Costa (USA) are Construction Global "
    "Ambassadors; Henrique applied and was accepted to lead the vertical at the Nucleo, with the "
    "Giga-projects & AI track. Timing: motto 'Megaprojects demand mega skills' (megaproject = over "
    "US$1bn, multi-year); AI is the lever for the three program pillars (safety, efficiency, "
    "sustainability), hub pmicp.us. Ask to the partner: co-curation plus access to the PMI-CP community "
    "for the founding cohort.",
    bordo=[
        "Henrique Diniz (Brazil) and Fabrício Costa (USA): Construction Global Ambassadors.",
        "Henrique was accepted to lead the vertical at the Núcleo, with the Giga-projects and AI track.",
        "Anchor proof: AI study on a typical megaproject claim or risk + webinar with the ambassadors."])

vertical("Vertical · PMO", "AI-augmented PMO: PMI-PMOCP",
    ["PMO leaders under pressure to prove value.",
     "Manual status reporting; scattered portfolio data.",
     "PMO seen as a cost, not as intelligence."],
    ["The PMO as an intelligence layer (portfolio analytics, forecasting).",
     "Automated status and reports straight from the data.",
     "AI-assisted portfolio prioritization."],
    "the hottest in the institutional model: PMI-PMOCP just launched (ISO, Feb/2026) and PMOGA absorbed "
    "by PMI. The community is eager for 'what AI changes in my PMO'.",
    "'Augmented PMO': how AI enters the 6 domains of PMI-PMOCP.",
    TEAL,
    "PMO vertical. Anchor credential PMI-PMOCP (successor to PMO-CP); partner PMO Global Alliance, now "
    "under PMI (pmoga.pmi.org). Maximum institutional alignment. Ask to the partner: a joint PMO + AI "
    "track and presence in the PMOGA community.")

vertical("Vertical · Agile", "Human judgment in the age of agents: PMI-ACP",
    ["Agilists rethinking the human role in the age of AI.",
     "Fear of 'AI versus agility'.",
     "Uncertainty about where human judgment adds value."],
    ["Jim Highsmith (co-author of the Agile Manifesto) reframes the question to 'what human leadership does not automate', and points to judgment as the most critical capability.",
     "Managing people and bots as a new agile competency.",
     "AI-accelerated delivery without losing the principles."],
    "Agile Alliance inside PMI + the PMP refresh (Jul/2026) emphasizing agile and hybrid; the community "
    "is repositioning.",
    "essay or debate 'human judgment in the age of agents', referencing the Manifesto.",
    BLUE,
    "Agile vertical. Anchor credential PMI-ACP; partner Agile Alliance (joined PMI in 2026). Thesis "
    "source: PMI AI Today, 'Reimagining Agility in an AI World'. SOURCE CAUTION: the popular 'five-hour "
    "sprints' phrase is NOT Highsmith's; use only the paraphrase of the framing (judgment, people and "
    "bots). Ask to the partner: a joint Agile Alliance x Nucleo session around the PMP refresh.")

vertical("Vertical · ESG / Green", "Sustainability is the #1 predictor of success: CSPP",
    ["Sustainability professionals in projects.",
     "Corporate ESG intent that does not turn into delivery.",
     "Hard measurement and reporting; scattered environmental and social data."],
    ["Sustainability measurement and reporting, aligned to the P5 Standard.",
     "AI turning ESG intent into traceable delivery.",
     "Analysis of environmental and social data at scale."],
    "the freshest: CSPP launched on 5 Jun/2026. PMI + GPM research (about 1,600 professionals, 35 "
    "countries) points to sustainability as the #1 predictor of success, ahead of methodology and governance.",
    "'AI + P5': how AI closes the execution gap. Ammo: 55% satisfaction (aligned) vs 33% (not); only 23% "
    "aligned today; 42-point confidence gap (85% sustainability vs 43% PMO).",
    GREEN,
    "ESG/Green vertical. Anchor credential CSPP (evolution of GPM-b, PMI + GPM, effective 5 Jun/2026); "
    "partner GPM Global. Numbers verified in the pitch kit: 55 vs 33; 23% aligned; 79% say it positions "
    "for the long term but only 41% integrate it; 42-point confidence gap (85% sustainability executives "
    "vs 43% PMO leaders). Media hook ready. Ask: a co-launch in the CSPP window.")

vertical("Vertical · Business", "AI in portfolio and strategy: PfMP, PgMP, PMI-PBA",
    ["Program and portfolio managers; business analysts.",
     "Strategic decisions under uncertainty.",
     "Portfolio prioritization; volatile requirements."],
    ["AI-assisted portfolio prioritization and scenarios.",
     "Augmented business analysis (requirements, stakeholders).",
     "Linking strategy and execution with data."],
    "PfMP, PgMP and PMI-PBA confirmed in the PMI registry. Three distinct sub-audiences: may become "
    "sub-verticals if demand justifies.",
    "portfolio prioritization case with AI.",
    RED,
    "Business/Program/Portfolio vertical. Anchor credentials PfMP (portfolio), PgMP (program), PMI-PBA "
    "(business analysis), confirmed in the registry. Partner to be defined; timing hook to verify. It is "
    "last in the activation order precisely because it depends on defining a partner first.")

# ---- 12..14 THE ASK (3 swappable variations) ----
def b_ask_board(s, top):
    add_box(s, 0.62, top, 12.1, 1.7, "For Mario Trentim and the PMI board", [
        "Strategic endorsement of the Núcleo as the horizontal seam across credential silos.",
        "Recognition of the Núcleo as a practical reading of PMI:Next and M.O.R.E.: the community that executes, at the human and research level, the silo integration PMI already does institutionally."],
        hcolor=PURPLE)
    add_box(s, 0.62, top+1.9, 12.1, 1.6, "Why it fits now", [
        "The move is already institutional (PMOGA, GPM, Agile Alliance joined PMI).",
        "The Núcleo gives a human face and research output to that strategy, with AI as the thread (PMIxAI, PMI Infinity)."],
        hcolor=TEAL)
content("The ask · PMI board", "What we ask the board", b_ask_board,
    "Variation 1 of 3 (swappable per audience). For the board: strategic endorsement and recognition of "
    "the Nucleo as a practical reading of PMI:Next / M.O.R.E. Anchor on the fact that silo integration is "
    "already institutional; the Nucleo is the human and research expression of it. Keep the other two "
    "variations hidden depending on the audience.")

def b_ask_chapter(s, top):
    add_box(s, 0.62, top, 12.1, 1.7, "For chapter presidents", [
        "Chapter adherence to the Núcleo IA & GP horizontal (light federation, without losing local identity).",
        "Nomination of protagonists, researchers and curators from the chapter for the verticals."],
        hcolor=PURPLE)
    add_box(s, 0.62, top+1.9, 12.1, 1.6, "What the chapter gains", [
        "Access to research output and to the Champion → CPMAI ladder for its members.",
        "Local validation of the AI-forward narrative, aligned with PMI's global strategy."],
        hcolor=TEAL)
content("The ask · Chapter presidents", "What we ask the chapters", b_ask_chapter,
    "Variation 2 of 3. For chapter presidents: chapter adherence, nomination of protagonists/researchers "
    "and validation of the narrative. Stress that it is a light federation, a research gain and a "
    "credential ladder for members, without losing local identity.")

def b_ask_partner(s, top):
    add_box(s, 0.62, top, 12.1, 1.7, "For vertical partners (GPM, Construction Ambassadors, PMOGA)", [
        "Co-curation of the vertical: content speaks the credential, the pain and the language of the community.",
        "Access to the credential community to form the founding cohort."],
        hcolor=PURPLE)
    add_box(s, 0.62, top+1.9, 12.1, 1.6, "The rule that protects the partner", [
        "The tribe produces; the vertical is a channel. No vertical owns the knowledge: it is a dock.",
        "Each vertical's pitch is refined with the partner before it goes public."],
        hcolor=TEAL)
content("The ask · Vertical partners", "What we ask the partners", b_ask_partner,
    "Variation 3 of 3. For vertical partners: co-curation plus access to the credential community for the "
    "founding cohort. Reinforce the anti-silo principle (production by the tribe, distribution by the "
    "vertical) as the guarantee that the partner does not lose control over its own knowledge.")

# ---- 15 NEXT STEPS ----
def b_next(s, top):
    add_table(s, [
        ["Activation order", "Why first", "Status"],
        ["1 · Construction", "Lead accepted (Henrique Diniz) + 2 Global Ambassadors (BR/US): the readiest to activate", "forming · lead in place"],
        ["2 · PMO", "PMI-PMOCP just launched + PMOGA absorbed: maximum institutional alignment", "forming · Cycle 4"],
        ["3 · ESG", "CSPP just launched (media window) + #1 predictor of success", "forming · Cycle 4"],
        ["4 · Agile", "Agile Alliance in PMI + PMP refresh (Jul)", "declared"],
        ["5 · Business", "define partner first", "declared"],
    ], top, widths=[1.8, 4.8, 1.9], fs=12)
    add_box(s, 0.62, top+2.95, 12.1, 1.1, "Be a protagonist", [
        "Each vertical enters with explicit status (forming, not vaporware). The CTA recruits founders, "
        "not consumers: a leadership program, aligned with M.O.R.E. and PMI:Next."],
        hcolor=ORANGE, hfs=15, fs=12.5)
content("Next steps", "Cycle 4: pilot verticals and a call for protagonists", b_next,
    "Closing. The activation order now opens by READINESS: Construction is #1 because it is the only one "
    "with an accepted lead (Henrique Diniz) and two Global Ambassadors (BR/US) already on board, with the "
    "Giga-projects & AI track. Then PMO (institutional alignment), ESG (CSPP media window, perishable), "
    "Agile (community repositioning), Business (define partner). Each vertical shows explicit status "
    "(forming), without faking activity. The call is 'Be a protagonist', not 'be a member': it recruits a "
    "founding cohort, with the M.O.R.E. and PMI:Next link that justifies calling it leadership. No "
    "hardcoding: the page reads the initiative status and renders the CTA.")

# delete template originals
for sldId in SRC: prs.slides._sldIdLst.remove(sldId)

# ===================== guards =====================
def _walk_text(shapes, where, fail):
    for sh in shapes:
        if sh.shape_type == 6:
            _walk_text(sh.shapes, where, fail); continue
        if getattr(sh, "has_table", False):
            for row in sh.table.rows:
                for cell in row.cells:
                    if "—" in cell.text_frame.text: fail(f"em-dash in {where}: table of {sh.name!r}")
            continue
        if sh.has_text_frame and "—" in sh.text_frame.text:
            fail(f"em-dash in {where}: shape {sh.name!r}")

def _assert_guards(p):
    def fail(msg): raise AssertionError(msg)
    sw, sh_ = p.slide_width, p.slide_height
    for i, sl in enumerate(p.slides, 1):
        _walk_text(sl.shapes, f"slide {i}", fail)
        if sl.has_notes_slide and "—" in (sl.notes_slide.notes_text_frame.text or ""):
            fail(f"em-dash in notes of slide {i}")
        for sp in sl.shapes:
            if sp.left is None: continue
            if sp.left + sp.width > sw + 9525 or sp.top + sp.height > sh_ + 9525:
                fail(f"slide {i}: shape {sp.name!r} overflows the canvas")
            if sp.name in BOILER: fail(f"slide {i}: boilerplate {sp.name!r} survived the wipe")
        dividers = [sp for sp in sl.shapes
                    if not sp.has_text_frame and sp.height < 18288 and sp.width > Inches(8)]
        for d in dividers:
            for (slide, l, t, w, h) in INJECTED:
                if slide is not sl: continue
                if t < d.top < t + h and not (l + w < d.left or d.left + d.width < l):
                    fail(f"slide {i}: divider {d.name!r} crosses injected content")
_assert_guards(prs)

prs.save(str(OUT))
print(f"saved {OUT} with {len(prs.slides._sldIdLst)} slides")

# ===================== render =====================
subprocess.run(["soffice", "--headless", "--convert-to", "pdf", "--outdir", str(BASE), str(OUT)],
               check=True, capture_output=True)
pdf = OUT.with_suffix(".pdf")
PREVIEW.mkdir(exist_ok=True)
for old in PREVIEW.glob("slide-*.png"): old.unlink()
shutil.copy2(pdf, PREVIEW / pdf.name)
subprocess.run(["pdftoppm", "-png", "-r", "100", str(pdf), str(PREVIEW / "slide")],
               check=True, capture_output=True)
n = len(list(PREVIEW.glob("slide-*.png")))
print(f"rendered {pdf.name} + {n} preview PNGs (from the final PDF)")
assert n == len(prs.slides._sldIdLst), "preview count != slide count"
