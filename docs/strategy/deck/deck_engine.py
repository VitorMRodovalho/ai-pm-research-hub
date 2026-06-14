#!/usr/bin/env python
# Shared engine for the Núcleo pitch deck — content-agnostic, language-agnostic.
# Clones the official PMI Events template's branded slides and injects by shape NAME,
# so the brand is preserved byte-for-byte and the file stays PowerPoint-editable.
# One Deck instance == one rendered deck. Build = pptx -> PDF -> preview PNGs in one run.
# Archetypes (this template): idx 10 = purple divider (cover); idx 17 = title+content.
# Guards fail the build: em-dash, canvas overflow, surviving boilerplate, divider-over-content.
import copy, subprocess, shutil
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn

# Official PMI palette (shared by every language edition)
PURPLE=RGBColor(0x46,0x1D,0xA3); PURPLE_DK=RGBColor(0x36,0x17,0x7B); LILAC=RGBColor(0xA9,0x8E,0xEC)
TEAL=RGBColor(0x44,0x78,0x9B); BLUE=RGBColor(0x6C,0xBE,0xDE); RED=RGBColor(0xB8,0x37,0x13)
ORANGE=RGBColor(0xE0,0x61,0x1F); GREEN=RGBColor(0x3E,0x8E,0x5A); TAN=RGBColor(0xB2,0x94,0x78)
LIGHT=RGBColor(0xF6,0xF4,0xEF); DARK=RGBColor(0x24,0x20,0x16); WHITE=RGBColor(0xFF,0xFF,0xFF)
GRAY=RGBColor(0x60,0x60,0x60)
FONT = "Aptos"
COVER, CONTENT, CONTENT_DARK = 10, 17, 18
BOILER = ("TextBox 9", "Group 10")          # light s17 body placeholder + footer how-to, removed per content slide
BOILER_DARK = ("TextBox 10", "Group 11")    # dark s18 equivalents (native PMI dark layout)
DARKBG = RGBColor(0x20, 0x0F, 0x3B)          # PMI dark brand background, as used by the template's own dark slides


def _iter(shapes):
    for sh in shapes:
        if sh.shape_type == 6: yield from _iter(sh.shapes)
        else: yield sh


class Deck:
    """One branded deck. SPECS (layout + content) drive it from the outside via the
    add_* / cover / content primitives; this class never hard-codes a slide."""

    def __init__(self, template, out, preview, hub, logo, dark=False):
        self.prs = Presentation(str(template))
        self._src = list(self.prs.slides._sldIdLst)
        self.slides = self.prs.slides
        self.out = Path(out); self.preview = Path(preview)
        self.hub = str(hub); self.logo = logo
        self.K = (self.prs.slide_width / 914400) / 13.333
        self.CT = 2.75 / self.K           # content top (design space) = below the template divider
        self._injected = []               # (slide, l, t, w, h) bboxes for the overlap guard
        self.dark = dark                  # dark theme: clone the template's native dark content layout
        self.body_default = LIGHT if dark else DARK

    # geometry helpers (design space 13.33 x 7.5 -> template EMU)
    def IN(self, v): return Inches(v * self.K)
    def FS(self, p): return Pt(round(p * self.K))

    def clone(self, src_idx):
        src = self.slides[src_idx]
        dst = self.slides.add_slide(src.slide_layout)
        for ph in list(dst.shapes):
            ph._element.getparent().remove(ph._element)
        rid_map = {}
        for rId, rel in src.part.rels.items():
            if rel.reltype.endswith(("slideLayout", "notesSlide")): continue
            rid_map[rId] = (dst.part.rels.get_or_add_ext_rel(rel.reltype, rel.target_ref)
                            if rel.is_external else dst.part.relate_to(rel.target_part, rel.reltype))
        spTree = dst.shapes._spTree
        for child in src.shapes._spTree.iterchildren():
            if child.tag.split('}')[-1] in ("nvGrpSpPr", "grpSpPr"): continue
            el = copy.deepcopy(child)
            for node in el.iter():
                for a in list(node.attrib):
                    if a.startswith('{http://schemas.openxmlformats.org/officeDocument/2006/relationships}'):
                        if node.get(a) in rid_map: node.set(a, rid_map[node.get(a)])
            spTree.append(el)
        # copy the slide-level background override (dark slides set <p:bg> at slide scope;
        # add_slide() only inherits the layout bg, so the dark fill is lost without this)
        src_cSld = src._element.find(qn('p:cSld'))
        dst_cSld = dst._element.find(qn('p:cSld'))
        src_bg = src_cSld.find(qn('p:bg')) if src_cSld is not None else None
        if src_bg is not None and dst_cSld is not None:
            old = dst_cSld.find(qn('p:bg'))
            if old is not None: dst_cSld.remove(old)
            dst_cSld.insert(0, copy.deepcopy(src_bg))
        return dst

    def named(self, slide, name): return [s for s in _iter(slide.shapes) if s.name == name]

    def delete(self, slide, *names):
        """Remove template shapes for real (clear() only empties text frames)."""
        for n in names:
            for s in list(slide.shapes):
                if s.name == n:
                    s._element.getparent().remove(s._element)

    def fix_year(self, slide):
        for sh in _iter(slide.shapes):
            if sh.has_text_frame and "2025" in sh.text_frame.text:
                for p in sh.text_frame.paragraphs:
                    for r in p.runs:
                        if "2025" in r.text: r.text = r.text.replace("2025", "2026")

    def inject_keep(self, shape, text):
        """Replace text but preserve the template run's size/bold/font/color."""
        tf = shape.text_frame
        runs = tf.paragraphs[0].runs
        size = bold = name = color = None
        if runs:
            f = runs[0].font; size, bold, name = f.size, f.bold, f.name
            try: color = f.color.rgb
            except Exception: color = None
        tf.clear()
        for i, ln in enumerate(text.split("\n")):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            r = p.add_run(); r.text = ln
            if size: r.font.size = size
            if bold is not None: r.font.bold = bold
            if name: r.font.name = name
            if color is not None:
                try: r.font.color.rgb = color
                except Exception: pass

    def notes(self, slide, text): slide.notes_slide.notes_text_frame.text = text

    def _track(self, slide, sp): self._injected.append((slide, sp.left, sp.top, sp.width, sp.height))

    def add_box(self, slide, left, top, w, h, header, lines, hcolor=PURPLE, fs=13, hfs=15, body=None):
        if body is None: body = self.body_default
        tb = slide.shapes.add_textbox(self.IN(left), self.IN(top), self.IN(w), self.IN(h))
        tf = tb.text_frame; tf.word_wrap = True
        first = True
        if header:
            p = tf.paragraphs[0]; r = p.add_run(); r.text = header
            r.font.bold = True; r.font.size = self.FS(hfs); r.font.name = FONT; r.font.color.rgb = hcolor
            p.space_after = Pt(8); first = False
        for ln in lines:
            p = tf.paragraphs[0] if first else tf.add_paragraph(); first = False
            r = p.add_run(); r.text = ln
            r.font.size = self.FS(fs); r.font.name = FONT; r.font.color.rgb = body
            p.space_after = Pt(6)
        self._track(slide, tb)
        return tb

    def add_table(self, slide, rows, top, left=0.62, width=12.1, fs=12, widths=None):
        nr, nc = len(rows), len(rows[0])
        gf = slide.shapes.add_table(nr, nc, self.IN(left), self.IN(top), self.IN(width), self.IN(0.4*nr))
        t = gf.table
        if widths:
            tot = sum(widths)
            for j, w in enumerate(widths): t.columns[j].width = self.IN(width*w/tot)
        for i, row in enumerate(rows):
            for j, val in enumerate(row):
                c = t.cell(i, j)
                for m in ("left","right"): setattr(c, f"margin_{m}", self.IN(0.08))
                for m in ("top","bottom"): setattr(c, f"margin_{m}", self.IN(0.03))
                c.fill.solid()
                if i == 0:
                    c.fill.fore_color.rgb = PURPLE; fc, fb = WHITE, True
                else:
                    c.fill.fore_color.rgb = WHITE if i % 2 else LIGHT; fc, fb = DARK, (j == 0)
                tf = c.text_frame; tf.word_wrap = True; tf.clear()
                p = tf.paragraphs[0]; r = p.add_run(); r.text = str(val)
                r.font.size = self.FS(fs); r.font.bold = fb; r.font.color.rgb = fc; r.font.name = FONT
        self._track(slide, gf)
        return gf

    def add_image(self, slide, top, width, left=None, path=None):
        if left is None: left = (13.333 - width) / 2
        pic = slide.shapes.add_picture(path or self.hub, self.IN(left), self.IN(top), width=self.IN(width))
        self._track(slide, pic)
        return pic

    def add_arrow(self, slide, left, top, color=LILAC, size=30):
        tb = slide.shapes.add_textbox(self.IN(left), self.IN(top), self.IN(0.9), self.IN(0.8))
        p = tb.text_frame.paragraphs[0]; r = p.add_run(); r.text = "→"
        r.font.size = self.FS(size); r.font.bold = True; r.font.name = FONT; r.font.color.rgb = color
        return tb

    def cover(self, title, subtitle, attribution, note=""):
        s = self.clone(COVER)
        ct = self.named(s, "TextBox 14"); cs = self.named(s, "TextBox 15")
        if ct: self.inject_keep(ct[0], title)
        if cs: self.inject_keep(cs[0], subtitle)
        self.fix_year(s)
        try:
            s.shapes.add_picture(self.logo, Inches(0.74), Inches(0.45), height=Inches(0.85))
        except Exception as e:
            print("logo skip:", e)
        attr = s.shapes.add_textbox(Inches(0.74), Inches(10.55), Inches(15.0), Inches(0.45))
        r = attr.text_frame.paragraphs[0].add_run(); r.text = attribution
        r.font.size = Pt(10); r.font.name = FONT; r.font.color.rgb = LILAC
        if note: self.notes(s, note)
        return s

    def content(self, eyebrow, title, builder, note=""):
        if self.dark:
            s = self.clone(CONTENT_DARK)
            title_name, eb_name, boiler = "TextBox 15", "TextBox 14", BOILER_DARK
        else:
            s = self.clone(CONTENT)
            title_name, eb_name, boiler = "TextBox 14", "TextBox 13", BOILER
        tb = self.named(s, title_name)
        if tb: self.inject_keep(tb[0], title)
        eb = self.named(s, eb_name)
        if eb: self.inject_keep(eb[0], eyebrow)
        self.delete(s, *boiler)
        self.fix_year(s)
        builder(s, self.CT)
        if note: self.notes(s, note)
        return s

    # ---- guards + save + render (last step of the build, same process) ----
    def _walk_text(self, shapes, where, fail):
        for sh in shapes:
            if sh.shape_type == 6:
                self._walk_text(sh.shapes, where, fail); continue
            if getattr(sh, "has_table", False):
                for row in sh.table.rows:
                    for cell in row.cells:
                        if "—" in cell.text_frame.text: fail(f"em-dash in {where}: table of {sh.name!r}")
                continue
            if sh.has_text_frame and "—" in sh.text_frame.text:
                fail(f"em-dash in {where}: shape {sh.name!r}")

    def _assert_guards(self):
        def fail(msg): raise AssertionError(msg)
        p = self.prs; sw, sh_ = p.slide_width, p.slide_height
        for i, sl in enumerate(p.slides, 1):
            self._walk_text(sl.shapes, f"slide {i}", fail)
            if sl.has_notes_slide and "—" in (sl.notes_slide.notes_text_frame.text or ""):
                fail(f"em-dash in notes of slide {i}")
            for sp in sl.shapes:
                if sp.left is None: continue
                if sp.left + sp.width > sw + 9525 or sp.top + sp.height > sh_ + 9525:
                    fail(f"slide {i}: shape {sp.name!r} overflows the canvas")
                if sp.name in (BOILER + BOILER_DARK): fail(f"slide {i}: boilerplate {sp.name!r} survived the wipe")
            dividers = [sp for sp in sl.shapes
                        if not sp.has_text_frame and sp.height < 18288 and sp.width > Inches(8)]
            for d in dividers:
                for (slide, l, t, w, h) in self._injected:
                    if slide is not sl: continue
                    if t < d.top < t + h and not (l + w < d.left or d.left + d.width < l):
                        fail(f"slide {i}: divider {d.name!r} crosses injected content")

    def _number_pages(self):
        """Replace the template page-number placeholder (literal '‹#›' or a slidenum field, which
        LibreOffice renders as '#') with a static run carrying the real number. Cover gets blank."""
        for i, sl in enumerate(self.prs.slides, 1):
            num = "" if i == 1 else str(i)
            for sh in _iter(sl.shapes):
                if not sh.has_text_frame: continue
                tf = sh.text_frame
                if tf.text.strip() not in ('‹#›', '#'): continue
                p0 = tf.paragraphs[0]
                size = name = None
                src_rpr = None
                r_el = p0._p.find(qn('a:r'))
                if r_el is not None: src_rpr = r_el.find(qn('a:rPr'))
                if src_rpr is None:
                    fld = p0._p.find(qn('a:fld'))
                    if fld is not None: src_rpr = fld.find(qn('a:rPr'))
                if src_rpr is not None:
                    if src_rpr.get('sz'): size = Pt(int(src_rpr.get('sz')) / 100)
                    latin = src_rpr.find(qn('a:latin'))
                    if latin is not None: name = latin.get('typeface')
                tf.clear()
                r = tf.paragraphs[0].add_run(); r.text = num
                if size: r.font.size = size
                if name: r.font.name = name
                # the template's own number color matches the dark bg (invisible) -> force theme color
                r.font.color.rgb = LIGHT if self.dark else DARK

    def finalize(self):
        """Delete template example slides, run guards, save, then render PDF + preview PNGs
        from the FINAL artifact in the same run (previews can never go stale)."""
        for sldId in self._src: self.prs.slides._sldIdLst.remove(sldId)
        self._number_pages()
        self._assert_guards()
        self.prs.save(str(self.out))
        n_slides = len(self.prs.slides._sldIdLst)
        print(f"saved {self.out.name} with {n_slides} slides")
        subprocess.run(["soffice", "--headless", "--convert-to", "pdf",
                        "--outdir", str(self.out.parent), str(self.out)],
                       check=True, capture_output=True)
        pdf = self.out.with_suffix(".pdf")
        self.preview.mkdir(exist_ok=True)
        for old in self.preview.glob("slide-*.png"): old.unlink()
        shutil.copy2(pdf, self.preview / pdf.name)
        subprocess.run(["pdftoppm", "-png", "-r", "100", str(pdf), str(self.preview / "slide")],
                       check=True, capture_output=True)
        n = len(list(self.preview.glob("slide-*.png")))
        print(f"rendered {pdf.name} + {n} preview PNGs (from the final PDF)")
        assert n == n_slides, "preview count != slide count"
