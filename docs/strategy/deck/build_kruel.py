#!/usr/bin/env python
# Build the COOPERATION pitch: Núcleo IA & GP × Grupo ALUN (interlocutor: Cristiano Kruel) — LIGHT theme.
#   ~/.venvs/pmo/bin/python gen_assets_kruel.py && ~/.venvs/pmo/bin/python build_kruel.py
# Reuses deck_engine.py (dark=False -> the PMI template's standard light content slide + numbered pages).
# Photos: assets/people/{vitor,fabricio}.jpg. Covers (drop-in, gitignored): assets/covers/*.{png,jpg};
# cert badges: assets/certs_proc/* (transparent, from gen_assets). Diagrams: assets/{strategy_flow,synergy}.png.
# Structure R3: light (brand-accurate), speaks the PMI Chapter Partnerships Framework. 11 slides, PT-BR.
from pathlib import Path
from PIL import Image
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Pt
from deck_engine import (Deck, PURPLE, TEAL, BLUE, GREEN, RED, ORANGE, LILAC, DARK, GRAY, LIGHT, FONT)
from deck_content_kruel import CONTENT

BASE = Path("/home/vitormrodovalho/projects/ai-pm-research-hub/docs/strategy/deck")
TEMPLATE = BASE / "build/pmi_template.pptx"
LOGO = "/home/vitormrodovalho/projects/_pmo/assets/pmi/brand/pmigo-logo-white.png"
PEOPLE = BASE / "assets/people"
COVERS = BASE / "assets/covers"
CERTS_PROC = BASE / "assets/certs_proc"
OUT = BASE / "Nucleo_IA_GP_Pitch_ALUN_Kruel.pptx"
PREVIEW = BASE / "preview_kruel"
HUB = BASE / "assets/hub_spoke.png"

# light-theme accents (PMI palette on the cream background)
A1, A2, A3 = PURPLE, TEAL, ORANGE


def panel(d, s, left, top, w, h, fill=None, line=None, lw=0.75):
    sp = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, d.IN(left), d.IN(top), d.IN(w), d.IN(h))
    if fill is not None:
        sp.fill.solid(); sp.fill.fore_color.rgb = fill
    else:
        sp.fill.background()
    if line is not None:
        sp.line.color.rgb = line; sp.line.width = Pt(lw)
    else:
        sp.line.fill.background()
    sp.shadow.inherit = False
    d._track(s, sp)
    return sp


def fitted_rect(left, top, boxw, boxh, path):
    with Image.open(path) as im:
        iw, ih = im.size
    ar = iw / ih
    if boxw / boxh > ar:
        h = boxh; w = h * ar
    else:
        w = boxw; h = w / ar
    return left + (boxw - w) / 2, top + (boxh - h) / 2, w, h


def place_fit(d, s, left, top, boxw, boxh, path):
    fl, ft, w, h = fitted_rect(left, top, boxw, boxh, path)
    pic = s.shapes.add_picture(path, d.IN(fl), d.IN(ft), width=d.IN(w), height=d.IN(h))
    d._track(s, pic)
    return pic


def img_slot(d, s, left, top, w, h, path, label, lfs=9, framed=False):
    if Path(path).exists():
        if framed:
            fl, ft, fw, fh = fitted_rect(left, top, w, h, str(path))
            panel(d, s, fl - 0.1, ft - 0.1, fw + 0.2, fh + 0.2, line=GRAY, lw=0.75)
        place_fit(d, s, left, top, w, h, str(path))
    else:
        d.add_box(s, left, top, w, h, "[ " + label + " ]", [], hcolor=GRAY, hfs=12, body=GRAY)
    if label:
        d.add_box(s, left, top + h + 0.03, w, 0.4, "", [label], fs=lfs, body=GRAY)


def cover_slot(d, s, left, top, w, h, slot, label, framed=False):
    found = None
    for ext in (".png", ".jpg", ".jpeg"):
        p = COVERS / f"{slot}{ext}"
        if p.exists():
            found = p; break
    img_slot(d, s, left, top, w, h, found or (COVERS / f"{slot}.png"), label, framed=framed)


def linked_line(d, s, left, top, w, h, segments, fs=11):
    tb = s.shapes.add_textbox(d.IN(left), d.IN(top), d.IN(w), d.IN(h))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]
    for text, url, color in segments:
        r = p.add_run(); r.text = text
        r.font.size = d.FS(fs); r.font.name = FONT; r.font.color.rgb = color
        if url:
            r.hyperlink.address = url
    d._track(s, tb)
    return tb


def two_col(d, s, top, h1, l1, h2, l2, c1=A1, c2=A2, bh=3.6, lw=5.95, rw=5.75, rl=6.95):
    d.add_box(s, 0.62, top, lw, bh, h1, l1, hcolor=c1)
    d.add_box(s, rl, top, rw, bh, h2, l2, hcolor=c2)


def caption(d, s, top, text, color=A1, w=12.1, fs=13):
    d.add_box(s, 0.62, top, w, 0.95, "", [text], hcolor=color, fs=fs)


def arrow3(d, s, left, top):
    tb = s.shapes.add_textbox(d.IN(left), d.IN(top), d.IN(0.6), d.IN(0.7))
    p = tb.text_frame.paragraphs[0]; r = p.add_run(); r.text = "→"
    r.font.size = d.FS(26); r.font.bold = True; r.font.name = FONT; r.font.color.rgb = LILAC


def compose(d, C):
    for k, v in C.items():
        t = v.get("title", "") if isinstance(v, dict) else ""
        if len(t) > 44:
            print(f"  WARN: title may wrap ({len(t)} chars): {t!r}")

    # 1 cover
    cov = C["cover"]
    d.cover(cov["title"], cov["sub"], cov["attr"], cov["note"])

    # 2 the fit (MAIVP wedge) + book cover
    f = C["fit"]
    def b_fit(s, top):
        top += 0.35
        d.add_box(s, 0.62, top, 4.25, 2.7, f["h1"], f["l1"], hcolor=A1)
        d.add_box(s, 5.0, top, 4.35, 2.7, f["h2"], f["l2"], hcolor=A2)
        cover_slot(d, s, 9.8, top - 0.1, 3.0, 3.4, "livro_kruel", "", framed=True)
        caption(d, s, top + 2.95, f["caption"], w=9.0)
    d.content(f["eyebrow"], f["title"], b_fit, f["note"])

    # 3 three fronts (menu) with the official PMI partnership TYPE under each
    fr = C["fronts"]
    def b_fronts(s, top):
        top += 0.45
        xs = (0.62, 4.74, 8.86); accents = (A1, A2, A3)
        for (head, ptype, lines), x, ac in zip(fr["cols"], xs, accents):
            d.add_box(s, x, top, 3.85, 0.5, head, [], hcolor=ac, hfs=15)
            d.add_box(s, x, top + 0.52, 3.85, 0.35, "", [ptype], fs=10, body=GRAY)
            d.add_box(s, x, top + 0.95, 3.85, 2.4, "", lines, fs=12)
        caption(d, s, top + 3.45, fr["caption"], color=A2)
    d.content(fr["eyebrow"], fr["title"], b_fronts, fr["note"])

    # 4 synergy / intersection (Venn hero + caption + the 6 PMI purposes)
    sy = C["synergy"]
    def b_syn(s, top):
        d.add_image(s, top + 0.0, 8.4, left=(13.333 - 8.4) / 2, path=str(BASE / "assets/synergy.png"))
        caption(d, s, top + 3.95, sy["caption"], color=A1, w=12.1)
        d.add_box(s, 0.62, top + 4.75, 12.1, 0.5, "", [sy["purposes"]], hcolor=A2, fs=11.5, body=GRAY)
    d.content(sy["eyebrow"], sy["title"], b_syn, sy["note"])

    # 5 who we are (people & culture first) + proof band (4 distinct lines)
    w = C["who"]
    def b_who(s, top):
        two_col(d, s, top, w["h1"], w["l1"], w["h2"], w["l2"], c1=A1, c2=A2, bh=2.8)
        d.add_box(s, 0.62, top + 2.95, 12.1, 1.75, w["proof_h"], w["proof"], hcolor=A3, hfs=14, fs=11)
    d.content(w["eyebrow"], w["title"], b_who, w["note"])

    # 6 the PMI (dedicated institution slide)
    pm = C["pmi"]
    def b_pmi(s, top):
        top += 0.35
        two_col(d, s, top, pm["h1"], pm["l1"], pm["h2"], pm["l2"], c1=A1, c2=A2, bh=2.9)
        caption(d, s, top + 3.05, pm["caption"], color=A3)
    d.content(pm["eyebrow"], pm["title"], b_pmi, pm["note"])

    # 7 ANSI authority + certifications (PT-BR legends)
    a = C["ansi"]
    def b_ansi(s, top):
        d.add_table(s, [a["head"]] + a["rows"], top + 0.05, left=0.62, width=7.1, widths=[2.1, 5.0], fs=10)
        cv = a["covers"][0]
        cover_slot(d, s, 8.2, top + 0.05, 4.3, 2.2, cv["slot"], cv["label"], framed=True)
        d.add_box(s, 0.62, top + 2.55, 12.1, 0.5, "", [a["caption"]], hcolor=A2, fs=11.5)
        n = len(a["certs"]); cw = 12.1 / n; bw = 1.05
        for i, c in enumerate(a["certs"]):
            cx = 0.62 + i * cw + (cw - bw) / 2
            proc = CERTS_PROC / (Path(c["file"]).stem + ".png")
            img = proc if proc.exists() else (BASE / "assets/inbox_r1" / c["file"])
            if img.exists():
                place_fit(d, s, cx, top + 3.0, bw, 1.0, str(img))
            else:
                d.add_box(s, cx, top + 3.0, bw, 1.0, "[ " + c["label"] + " ]", [], hcolor=GRAY, hfs=10, body=GRAY)
            hc = A3 if c.get("hero") else A1
            d.add_box(s, 0.62 + i * cw + 0.05, top + 4.05, cw - 0.1, 0.75, c["label"], [c["desc"]], hcolor=hc, hfs=10.5, fs=9)
    d.content(a["eyebrow"], a["title"], b_ansi, a["note"])

    # 8 why now (market problem, de-denigrated)
    pr = C["whynow"]
    def b_why(s, top):
        d.add_image(s, top + 0.15, 7.2, left=0.4, path=str(BASE / "assets/strategy_flow.png"))
        d.add_box(s, 0.5, top + 2.4, 7.1, 2.1, "", [pr["caption"]], hcolor=A3, fs=12.5)
        cv = pr["covers"]
        cover_slot(d, s, 8.1, top + 0.2, 4.2, 1.85, cv[0]["slot"], cv[0]["label"], framed=True)
        cover_slot(d, s, 8.1, top + 2.55, 4.2, 1.85, cv[1]["slot"], cv[1]["label"], framed=True)
    d.content(pr["eyebrow"], pr["title"], b_why, pr["note"])

    # 9 how we formalize (governance path, authority-framed)
    fo = C["formalize"]
    def b_form(s, top):
        top += 0.4; h = 3.0
        (h1, l1), (h2, l2), (h3, l3) = fo["steps"]
        d.add_box(s, 0.62, top, 3.5, h, h1, l1, hcolor=A1)
        arrow3(d, s, 4.25, top + 1.0)
        d.add_box(s, 4.85, top, 3.5, h, h2, l2, hcolor=A2)
        arrow3(d, s, 8.5, top + 1.0)
        d.add_box(s, 9.1, top, 3.6, h, h3, l3, hcolor=A3)
        caption(d, s, top + 3.2, fo["caption"], color=A1, w=12.1)
    d.content(fo["eyebrow"], fo["title"], b_form, fo["note"])

    # 10 who leads (photos + clickable contacts)
    mg = C["managers"]
    def b_mgr(s, top):
        pw, txw = 2.0, 3.6
        for left, person in ((0.7, mg["people"][0]), (6.9, mg["people"][1])):
            photo = PEOPLE / person["photo"]
            if photo.exists():
                d.add_image(s, top + 0.05, pw, left=left, path=str(photo))
            else:
                d.add_box(s, left, top + 0.05, pw, pw, "[ FOTO ]", [person["photo"]],
                          hcolor=GRAY, hfs=14, fs=10, body=GRAY)
            tx = left + pw + 0.3
            d.add_box(s, tx, top + 0.05, txw, 1.9, person["name"],
                      [person["role"], person["sub"]], hcolor=A1, hfs=15, fs=11)
            linked_line(d, s, tx, top + 2.05, txw, 0.4,
                        [("LinkedIn: " + person["linkedin"], "https://" + person["linkedin"], A2)], fs=11)
            linked_line(d, s, tx, top + 2.45, txw, 0.4,
                        [("Tel: ", None, GRAY), (person["phone"], "tel:" + person["phone"].replace(" ", ""), A2)], fs=11)
        segs = []
        for i, (lbl, shown, url) in enumerate(mg["links"]):
            if i: segs.append(("    ·    ", None, GRAY))
            segs.append((lbl + ": ", None, GRAY)); segs.append((shown, url, A2))
        linked_line(d, s, 0.62, top + 4.05, 12.1, 0.5, segs, fs=11)
    d.content(mg["eyebrow"], mg["title"], b_mgr, mg["note"])

    # 11 the ask (ambitious close)
    ak = C["ask"]
    def b_ask(s, top):
        top += 0.5
        two_col(d, s, top, ak["h1"], ak["l1"], ak["h2"], ak["l2"], c1=A1, c2=A2, bh=2.8)
        d.add_box(s, 0.62, top + 3.0, 12.1, 0.95, ak["cta_h"], [ak["cta"]], hcolor=A3, hfs=15, fs=13)
    d.content(ak["eyebrow"], ak["title"], b_ask, ak["note"])


def main():
    d = Deck(TEMPLATE, OUT, PREVIEW, HUB, LOGO, dark=False)
    compose(d, CONTENT["pt"])
    d.finalize()


if __name__ == "__main__":
    main()
