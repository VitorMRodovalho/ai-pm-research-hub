#!/usr/bin/env python
# Núcleo IA & GP — executive pitch deck (board PMI / chapter presidents / vertical partners).
# Brand: PMIGO via the OFFICIAL PMI Events template (clone + inject by shape name).
# Engine reused from pmigo-plataforma/decks/build/build_deck.py (proven clone/inject/guards/QA).
# Content is a RENDER of docs/strategy/{verticals_x_quadrants_model, vertical_pitch_kit,
# cycle4_landing_value_prop}.md + deck_outline.md. No research here, only transposition.
# Archetypes: idx 10 (purple divider) = cover; idx 17 (title+content) = content.
# Build = pptx -> PDF -> preview PNGs in ONE run (previews can never go stale).
# Guards: em-dash, slide-bounds, leftover-boilerplate, divider-overlap. No em-dash anywhere.
import copy, subprocess, shutil
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

BASE = Path("/home/vitormrodovalho/projects/ai-pm-research-hub/docs/strategy/deck")
TEMPLATE = BASE / "build/pmi_template.pptx"
OUT = BASE / "Nucleo_IA_GP_Pitch_Executivo.pptx"
PREVIEW = BASE / "preview"
HUB = str(BASE / "assets/hub_spoke.png")
LOGO = "/home/vitormrodovalho/projects/19SGPL-PMIGO/MKT/Pending/pmi_chp_logo_goias_brazil_hrz_wht.png"

# Official PMI palette
PURPLE=RGBColor(0x46,0x1D,0xA3); PURPLE_DK=RGBColor(0x36,0x17,0x7B); LILAC=RGBColor(0xA9,0x8E,0xEC)
TEAL=RGBColor(0x44,0x78,0x9B); BLUE=RGBColor(0x6C,0xBE,0xDE); RED=RGBColor(0xB8,0x37,0x13)
ORANGE=RGBColor(0xE0,0x61,0x1F); GREEN=RGBColor(0x3E,0x8E,0x5A); TAN=RGBColor(0xB2,0x94,0x78)
LIGHT=RGBColor(0xF6,0xF4,0xEF); DARK=RGBColor(0x24,0x20,0x16); WHITE=RGBColor(0xFF,0xFF,0xFF)
GRAY=RGBColor(0x60,0x60,0x60)
FONT="Aptos"
COVER, CONTENT = 10, 17

prs = Presentation(str(TEMPLATE))
SRC = list(prs.slides._sldIdLst)
slides = prs.slides
SW_IN = prs.slide_width / 914400; SH_IN = prs.slide_height / 914400
K = SW_IN / 13.333
CT = 2.75 / K            # content top (design space) = below the divider line
INJECTED = []

def IN(v): return Inches(v * K)
def FS(p): return Pt(round(p * K))

def _iter(shapes):
    for sh in shapes:
        if sh.shape_type == 6: yield from _iter(sh.shapes)
        else: yield sh

def clone(src_idx):
    src = slides[src_idx]
    dst = slides.add_slide(src.slide_layout)
    for ph in list(dst.shapes):
        ph._element.getparent().remove(ph._element)
    rid_map = {}
    for rId, rel in src.part.rels.items():
        if rel.reltype.endswith(("slideLayout", "notesSlide")): continue
        rid_map[rId] = (dst.part.rels.get_or_add_ext_rel(rel.reltype, rel.target_ref)
                        if rel.is_external else dst.part.relate_to(rel.target_part, rel.reltype))
    spTree = dst.shapes._spTree
    for child in src.shapes._spTree.iterchildren():
        if child.tag.split('}')[-1] in ("nvGrpSpPr", "grpSpPr"): continue
        el = copy.deepcopy(child)
        for node in el.iter():
            for a in list(node.attrib):
                if a.startswith('{http://schemas.openxmlformats.org/officeDocument/2006/relationships}'):
                    if node.get(a) in rid_map: node.set(a, rid_map[node.get(a)])
        spTree.append(el)
    return dst

def named(slide, name): return [s for s in _iter(slide.shapes) if s.name == name]

def delete(slide, *names):
    for n in names:
        for s in list(slide.shapes):
            if s.name == n:
                s._element.getparent().remove(s._element)

def fix_year(slide):
    for sh in _iter(slide.shapes):
        if sh.has_text_frame and "2025" in sh.text_frame.text:
            for p in sh.text_frame.paragraphs:
                for r in p.runs:
                    if "2025" in r.text: r.text = r.text.replace("2025", "2026")

def inject_keep(shape, text):
    tf = shape.text_frame
    runs = tf.paragraphs[0].runs
    size = bold = name = color = None
    if runs:
        f = runs[0].font; size, bold, name = f.size, f.bold, f.name
        try: color = f.color.rgb
        except Exception: color = None
    tf.clear()
    for i, ln in enumerate(text.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        r = p.add_run(); r.text = ln
        if size: r.font.size = size
        if bold is not None: r.font.bold = bold
        if name: r.font.name = name
        if color is not None:
            try: r.font.color.rgb = color
            except Exception: pass

def notes(slide, text): slide.notes_slide.notes_text_frame.text = text

def _track(slide, sp): INJECTED.append((slide, sp.left, sp.top, sp.width, sp.height))

def add_box(slide, left, top, w, h, header, lines, hcolor=PURPLE, fs=13, hfs=15, body=DARK):
    tb = slide.shapes.add_textbox(IN(left), IN(top), IN(w), IN(h))
    tf = tb.text_frame; tf.word_wrap = True
    first = True
    if header:
        p = tf.paragraphs[0]; r = p.add_run(); r.text = header
        r.font.bold = True; r.font.size = FS(hfs); r.font.name = FONT; r.font.color.rgb = hcolor
        p.space_after = Pt(8); first = False
    for ln in lines:
        p = tf.paragraphs[0] if first else tf.add_paragraph(); first = False
        r = p.add_run(); r.text = ln
        r.font.size = FS(fs); r.font.name = FONT; r.font.color.rgb = body
        p.space_after = Pt(6)
    _track(slide, tb)
    return tb

def add_table(slide, rows, top, left=0.62, width=12.1, fs=12, widths=None):
    nr, nc = len(rows), len(rows[0])
    gf = slide.shapes.add_table(nr, nc, IN(left), IN(top), IN(width), IN(0.4*nr))
    t = gf.table
    if widths:
        tot = sum(widths)
        for j, w in enumerate(widths): t.columns[j].width = IN(width*w/tot)
    for i, row in enumerate(rows):
        for j, val in enumerate(row):
            c = t.cell(i, j)
            for m in ("left","right"): setattr(c, f"margin_{m}", IN(0.08))
            for m in ("top","bottom"): setattr(c, f"margin_{m}", IN(0.03))
            c.fill.solid()
            if i == 0:
                c.fill.fore_color.rgb = PURPLE; fc, fb = WHITE, True
            else:
                c.fill.fore_color.rgb = WHITE if i % 2 else LIGHT; fc, fb = DARK, (j == 0)
            tf = c.text_frame; tf.word_wrap = True; tf.clear()
            p = tf.paragraphs[0]; r = p.add_run(); r.text = str(val)
            r.font.size = FS(fs); r.font.bold = fb; r.font.color.rgb = fc; r.font.name = FONT
    _track(slide, gf)
    return gf

def add_image(slide, path, top, width, left=None):
    if left is None: left = (13.333 - width) / 2
    pic = slide.shapes.add_picture(path, IN(left), IN(top), width=IN(width))
    _track(slide, pic)
    return pic

def add_arrow(slide, left, top, color=PURPLE, size=30):
    tb = slide.shapes.add_textbox(IN(left), IN(top), IN(0.9), IN(0.8))
    p = tb.text_frame.paragraphs[0]; r = p.add_run(); r.text = "→"
    r.font.size = FS(size); r.font.bold = True; r.font.name = FONT; r.font.color.rgb = color
    return tb

BOILER = ("TextBox 9", "Group 10")

def content(eyebrow, title, builder, note=""):
    s = clone(CONTENT)
    tb = named(s, "TextBox 14")
    if tb: inject_keep(tb[0], title)
    eb = named(s, "TextBox 13")
    if eb: inject_keep(eb[0], eyebrow)
    delete(s, *BOILER)
    fix_year(s)
    builder(s, CT)
    if note: notes(s, note)
    return s

# ===================== SPECS (design space 13.33 x 7.5; sem em-dash) =====================

# ---- 1 COVER ----
cov = clone(COVER)
ct = named(cov, "TextBox 14"); cs = named(cov, "TextBox 15")
if ct: inject_keep(ct[0], "Núcleo IA & GP")
if cs: inject_keep(cs[0], "A costura entre os silos de credencial do PMI  ·  "
                          "Núcleo de Estudos e Pesquisa em IA & GP  ·  Junho de 2026")
fix_year(cov)
try:
    cov.shapes.add_picture(LOGO, Inches(0.74), Inches(0.45), height=Inches(0.85))
except Exception as e: print("logo skip:", e)
attr = cov.shapes.add_textbox(Inches(0.74), Inches(10.55), Inches(15.0), Inches(0.45))
ap = attr.text_frame.paragraphs[0]; ar = ap.add_run()
ar.text = ("PMI, o logotipo PMI, PMI-CP, PMI-PMOCP, PMI-ACP, CSPP, CPMAI, PfMP, PgMP e PMI-PBA "
           "são marcas do Project Management Institute, Inc. Uso conforme as diretrizes de marca para capítulos.")
ar.font.size = Pt(10); ar.font.name = FONT; ar.font.color.rgb = LILAC
notes(cov, "Abertura. O Núcleo IA & GP é a costura entre os silos de credencial do PMI: uma "
           "comunidade voluntária que junta gente boa de toda credencial em pesquisa, desenvolvimento "
           "e networking, com a IA como o fio. Este deck mostra a tese (a horizontal que o PMI nao tem "
           "por desenho), o modelo de tres eixos, as verticais e o pedido por audiencia. Publico: board "
           "do PMI, presidentes de capitulo e parceiros de vertical. Marca PMIGO sob as diretrizes de "
           "capitulo; frase de atribuicao no rodape.")

# ---- 2 O PROBLEMA (modelo §1) ----
def b_prob(s, top):
    add_box(s, 0.62, top, 5.95, 4.3, "O PMI vive em silos de credencial", [
        "Construção, PMO, Ágil, Sustentabilidade, Negócio.",
        "Cada comunidade na sua trilha; raramente conversam.",
        "Toda parceria (GPM, Construction Ambassadors, PMOGA) vira caso especial e isolado."],
        hcolor=PURPLE)
    add_box(s, 6.95, top, 5.75, 4.3, "A IA é transversal a todas", [
        "Nenhuma credencial escapa da IA: ela atravessa todos os silos.",
        "Falta no ecossistema uma horizontal que costure as fronteiras.",
        "Esse é o papel do Núcleo: voluntário para voluntário, um só propósito, gente boa através das credenciais."],
        hcolor=ORANGE)
content("O problema", "O ativo que falta no PMI é uma horizontal", b_prob,
    "O problema nao e falta de comunidades, e falta de uma costura entre elas. O PMI e organizado por "
    "silos de credencial; a IA e transversal a todos. Sem uma horizontal explicita, cada parceria vira "
    "caso especial. O Nucleo e essa horizontal por desenho: comunidade voluntaria, um proposito, "
    "juntando gente boa atraves das fronteiras de credencial, com a IA como o fio.")

# ---- 3 O FIO + VENTO A FAVOR (modelo §1.1) -- slide-matador ----
def b_tail(s, top):
    add_table(s, [
        ["Silo", "Movimento do próprio PMI (verificado, jun/2026)"],
        ["PMO", "PMI adquiriu o PMO-CP da PMOGA (2023) e relançou como PMI-PMOCP, ISO-accredited (fev/2026); PMOGA hoje em pmoga.pmi.org"],
        ["ESG / Verde", "GPM-b evolui para CSPP, co-branded PMI + GPM, alinhado ao Standard P5 (efetivo 5 jun/2026)"],
        ["Ágil", "Agile Alliance entrou no PMI (2026)"],
    ], top, widths=[1.4, 6.8], fs=13)
    add_box(s, 0.62, top+2.95, 12.1, 1.1,
        "O Núcleo é a expressão humana e de pesquisa dessa mesma integração. A IA é o fio.", [],
        hcolor=PURPLE, hfs=15)
content("O vento a favor", "Não remamos contra a maré: nomeamos a maré", b_tail,
    "Este e o slide-matador para o board. A tese de hub integrador nao rema contra a mare "
    "institucional, ela nomeia a mare: o proprio PMI esta absorvendo comunidades de credencial para "
    "dentro de si (PMO-CP da PMOGA em 2023; GPM-b vira CSPP em jun/2026; Agile Alliance entra em 2026). "
    "O Nucleo e a expressao humana e de pesquisa dessa integracao, com a IA como fio. Leitura pratica de "
    "PMI:Next e M.O.R.E., nao mais uma iniciativa.")

# ---- 4 O MODELO (modelo §3 + landing §5a) -- hub-and-spoke ----
def b_model(s, top):
    add_image(s, HUB, top-0.15, 9.7)
    add_box(s, 0.62, top+4.35, 12.1, 1.0, "", [
        "Três eixos ortogonais: Quadrante (o quê) × Tribo (quem produz) × Vertical (pra quem aterrissa). "
        "A escada comum a toda vertical: PMIxAI Champion (aberto) → Grupo de Estudos CPMAI → PMI-CPMAI.",
        "A vertical é doca, nunca silo: a tribo produz, a vertical distribui. Ninguém é dono do conhecimento."],
        hcolor=PURPLE, fs=12, hfs=13, body=DARK)
content("O modelo", "A IA no centro, as comunidades como raios", b_model,
    "O visual e o pitch. Centro: Nucleo + IA, a costura. Raios: as verticais, cada uma uma comunidade de "
    "credencial do PMI. Anel: a escada Champion > CPMAI, espinha comum a todos os raios. O modelo tem "
    "tres eixos ortogonais: quadrante (tipo de conhecimento, taxonomia propria do Nucleo), tribo (quem "
    "produz, Eixo A) e vertical (pra quem distribui, Eixo B). A vertical nao duplica o quadrante: define "
    "publico e empacotamento. Principio anti-silo: producao e da tribo, distribuicao e da vertical; "
    "nenhuma vertical e dona de conhecimento, e doca.")

# ---- 5 A ESCADA (modelo §4) ----
def b_ladder(s, top):
    y = top + 0.35; h = 3.0
    add_box(s, 0.62, y, 3.5, h, "1 · PMIxAI Champion", [
        "Reconhecimento e badge, aberto e leve.",
        "Porta de entrada (Eixo B).",
        "Já é primitivo da plataforma (award_champion)."], hcolor=BLUE)
    add_arrow(s, 4.25, y+1.0, color=LILAC)
    add_box(s, 4.85, y, 3.5, h, "2 · Grupo de Estudos CPMAI", [
        "Preparatório, em tribo (Eixo A).",
        "A ponte entre o aberto e a credencial.",
        "Já modelado como iniciativa (cpmai_*)."], hcolor=TEAL)
    add_arrow(s, 8.5, y+1.0, color=LILAC)
    add_box(s, 9.1, y, 3.6, h, "3 · PMI-CPMAI", [
        "A credencial-costura.",
        "Faz interface com TODAS as verticais.",
        "O ponto comum entre comunidades que não se falariam."], hcolor=PURPLE)
content("A escada comum", "A costura de credencial é uma escada, não um ponto", b_ladder,
    "A costura tem uma escada que toda vertical atravessa. PMIxAI Champion (aberto, gamificado, ja existe "
    "como award_champion) prepara o terreno; o Grupo de Estudos CPMAI (Eixo A, preparatorio, ja modelado) "
    "faz a ponte; o PMI-CPMAI e a credencial-costura que faz interface com todas as verticais. Isomorfia "
    "central: no ecossistema PMI os silos sao credenciais e a costura e a IA (CPMAI); dentro do Nucleo a "
    "costura e o mesmo mecanismo, com a escada Champion > CPMAI como espinha compartilhada.")

# ---- 6 COBERTURA & ALCANCE (landing §5b) ----
def b_reach(s, top):
    add_box(s, 0.62, top, 5.95, 4.0, "Brasil + LatAm em primeiro plano", [
        "Membros ativos (pesquisadores, líderes, curadores) já em vários estados do Brasil.",
        "15 capítulos no Brasil; diversidade intra-Brasil é a prioridade.",
        "Enquadrar a LatAm é o recado de expansão por geografia, sem precisar declarar em texto."],
        hcolor=PURPLE)
    add_box(s, 6.95, top, 5.75, 4.0, "Presença internacional = networking", [
        "O Núcleo já reúne gente em Brasil, Portugal, Itália e EUA.",
        "Não é heatmap mundial esparso: é ativo de acesso a networking, nomeado por país.",
        "Embaixadores são poucos (cerca de 4); a comunidade ativa é maior que isso."],
        hcolor=TEAL)
    add_box(s, 0.62, top+4.15, 12.1, 0.55,
        "Governança-first (LGPD): agregação por capítulo, estado ou país. Zero PII; pin individual só com opt-in.",
        [], hcolor=ORANGE, hfs=12)
content("Cobertura & alcance", "Quebramos o silo geográfico, não só o de credencial", b_reach,
    "Calibracao de jun/2026: o Nucleo nao e so embaixadores (sao cerca de 4); membros ativos ja estao em "
    "varios estados do Brasil mais Portugal, Italia e EUA. A diversidade prioritaria e intra-Brasil, mas a "
    "presenca internacional e ativo de networking e deve aparecer, nomeada por pais, nao como borrao "
    "mundial. Em pagina publica tudo agrega por capitulo/estado/pais, zero PII, pin de pessoa so com "
    "opt-in registrado. O mapa do PMAIrevolution pode ser reaproveitado, re-projetado para Brasil/LatAm.")

# ---- 7..11 VERTICAIS (vertical_pitch_kit) ----
def vertical(eyebrow, title, dor, teses, timing, prova, hc, note, bordo=None):
    def b(s, top):
        bh = 3.4 if bordo else 4.1
        add_box(s, 0.62, top, 5.95, bh, "Público & dor", dor, hcolor=hc)
        add_box(s, 6.95, top, 5.75, bh, "Tese de IA", teses, hcolor=PURPLE)
        if bordo:
            add_box(s, 0.62, top+bh+0.15, 12.1, 1.55, "Já a bordo: liderança definida",
                    bordo, hcolor=ORANGE, hfs=14, fs=12.5)
        else:
            add_box(s, 0.62, top+4.2, 12.1, 0.6, "", [f"Timing: {timing}", f"Prova-âncora: {prova}"],
                    fs=11.5, body=DARK)
    return content(eyebrow, title, b, note)

vertical("Vertical · Construção", "Giga-projetos: PMI-CP e a IA nos megaprojetos",
    ["Líderes de megaprojetos e infraestrutura (A/E/C).",
     "Projetos longos, documentação massiva, risco de prazo e custo.",
     "Contratos e RFIs intermináveis; dados fragmentados entre stakeholders."],
    ["Análise de risco e cronograma em megaprojetos.",
     "Leitura e sumarização de contratos, RFIs e submittals em escala.",
     "Gêmeos digitais e dados de campo como insumo de decisão."],
    "mote do PMI 'Megaprojects demand mega skills' (megaprojeto = mais de US$1bi, multi-ano); "
    "Construction Ambassadors advoga segurança, eficiência e sustentabilidade (hub pmicp.us).",
    "estudo 'IA aplicada a um pleito ou risco típico de megaprojeto' + webinar com um embaixador.",
    ORANGE,
    "Vertical Construcao. Credencial-ancora PMI-CP; parceiro Global Construction Ambassadors. JA TEM "
    "protagonistas: Henrique Diniz (Brasil) e Fabricio Costa (EUA) sao Construction Global Ambassadors; "
    "Henrique se inscreveu e foi aceito para liderar a vertical no Nucleo, com a trilha Giga Projetos e "
    "IA. Timing: mote 'Megaprojects demand mega skills' (megaprojeto = mais de US$1bi, multi-ano); IA e a "
    "alavanca dos tres pilares (seguranca, eficiencia, sustentabilidade), hub pmicp.us. Pedido ao "
    "parceiro: co-curadoria mais acesso a comunidade PMI-CP para a coorte fundadora.",
    bordo=[
        "Henrique Diniz (Brasil) e Fabrício Costa (EUA): Construction Global Ambassadors.",
        "Henrique foi aceito para liderar a vertical no Núcleo, com a trilha Giga Projetos e IA.",
        "Prova-âncora: estudo de IA em pleito ou risco de megaprojeto + webinar com os embaixadores."])

vertical("Vertical · PMO", "PMO aumentado por IA: PMI-PMOCP",
    ["Líderes de PMO sob pressão de provar valor.",
     "Status reporting manual; dados de portfólio dispersos.",
     "PMO visto como custo, não como inteligência."],
    ["O PMO como camada de inteligência (analytics de portfólio, previsão).",
     "Status e relatórios automatizados a partir dos dados.",
     "Priorização de portfólio assistida por IA."],
    "o mais quente do modelo institucional: PMI-PMOCP recém-lançada (ISO, fev/2026) e a PMOGA "
    "absorvida pelo PMI. Comunidade ávida por 'o que a IA muda no meu PMO'.",
    "'PMO Aumentado': como a IA entra nos 6 domínios do PMI-PMOCP.",
    TEAL,
    "Vertical PMO. Credencial-ancora PMI-PMOCP (sucessora do PMO-CP); parceiro PMO Global Alliance, hoje "
    "sob o PMI (pmoga.pmi.org). Maximo alinhamento institucional: e a primeira da ordem de ativacao. "
    "Pedido ao parceiro: trilha conjunta PMO + IA e presenca na comunidade PMOGA.")

vertical("Vertical · Ágil", "Julgamento humano na era dos agentes: PMI-ACP",
    ["Agilistas repensando o papel humano na era da IA.",
     "Medo de 'IA contra agilidade'.",
     "Incerteza sobre onde o julgamento humano agrega."],
    ["Jim Highsmith (co-autor do Manifesto Ágil) reposiciona a pergunta para 'que liderança humana não se automatiza', e aponta o julgamento como a capacidade mais crítica.",
     "Gerir pessoas e bots como nova competência ágil.",
     "Entrega acelerada por IA sem perder os princípios."],
    "Agile Alliance dentro do PMI + refresh do PMP (jul/2026) enfatizando ágil e híbrido; "
    "a comunidade está se reposicionando.",
    "ensaio ou debate 'julgamento humano na era dos agentes', com referência ao Manifesto.",
    BLUE,
    "Vertical Agil. Credencial-ancora PMI-ACP; parceiro Agile Alliance (entrou no PMI em 2026). Fonte da "
    "tese: PMI AI Today, 'Reimagining Agility in an AI World'. ATENCAO de fonte: a frase popular de "
    "'sprints de 5h' NAO e do Highsmith; usar apenas a parafrase do enquadramento (julgamento, pessoas e "
    "bots). Pedido ao parceiro: sessao conjunta Agile Alliance x Nucleo na virada do PMP.")

vertical("Vertical · ESG / Verde", "Sustentabilidade é o #1 preditor de sucesso: CSPP",
    ["Profissionais de sustentabilidade em projetos.",
     "Intenção corporativa de ESG que não vira entrega.",
     "Medição e reporte difíceis; dados ambientais e sociais dispersos."],
    ["Medição e reporte de sustentabilidade, alinhados ao Standard P5.",
     "IA tornando intenção ESG em entrega rastreável.",
     "Análise de dados ambientais e sociais em escala."],
    "o mais fresco: CSPP saiu em 5 jun/2026. Pesquisa PMI + GPM (cerca de 1.600 profissionais, 35 "
    "países) aponta sustentabilidade como o #1 preditor de sucesso, à frente de metodologia e governança.",
    "'IA + P5': como a IA fecha o gap de execução. Munição: 55% de satisfação (alinhados) vs 33% (não); "
    "só 23% alinhados hoje; gap de confiança de 42 pontos (85% sust. vs 43% PMO).",
    GREEN,
    "Vertical ESG/Verde. Credencial-ancora CSPP (evolucao do GPM-b, PMI + GPM, efetiva 5 jun/2026); "
    "parceiro GPM Global. Numeros conferidos no pitch kit: 55 vs 33; 23% alinhados; 79% dizem que "
    "posiciona para o longo prazo mas so 41% integram; gap de confianca de 42 pontos (85% executivos de "
    "sustentabilidade vs 43% lideres de PMO). Gancho de midia pronto. Pedido: co-lancamento na janela CSPP.")

vertical("Vertical · Negócio", "IA em portfólio e estratégia: PfMP, PgMP, PMI-PBA",
    ["Gestores de programa e portfólio; analistas de negócio.",
     "Decisão estratégica sob incerteza.",
     "Priorização de portfólio; requisitos voláteis."],
    ["Priorização de portfólio e cenários assistidos por IA.",
     "Análise de negócio aumentada (requisitos, stakeholders).",
     "Ligação estratégia e execução com dados."],
    "PfMP, PgMP e PMI-PBA confirmadas no registry PMI. São 3 sub-públicos distintos: podem virar "
    "sub-verticais se a demanda justificar.",
    "caso de priorização de portfólio com IA.",
    RED,
    "Vertical Negocio/Programa/Portfolio. Credenciais-ancora PfMP (portfolio), PgMP (programa), PMI-PBA "
    "(analise de negocio), confirmadas no registry. Parceiro a definir; gancho de timing a verificar. "
    "E a ultima da ordem de ativacao justamente por depender de definir parceiro primeiro.")

# ---- 12 O PEDIDO (3 variações trocáveis) ----
def b_ask_board(s, top):
    add_box(s, 0.62, top, 12.1, 1.7, "Para Mario Trentim e o board do PMI", [
        "Endorsement estratégico do Núcleo como horizontal-costura entre os silos de credencial.",
        "Reconhecimento do Núcleo como leitura prática de PMI:Next e M.O.R.E.: a comunidade que executa, no nível humano e de pesquisa, a integração de silos que o PMI já faz institucionalmente."],
        hcolor=PURPLE)
    add_box(s, 0.62, top+1.9, 12.1, 1.6, "Por que cabe agora", [
        "O movimento já é institucional (PMOGA, GPM, Agile Alliance entraram para dentro do PMI).",
        "O Núcleo dá rosto humano e produção de pesquisa a essa estratégia, com a IA como fio (PMIxAI, PMI Infinity)."],
        hcolor=TEAL)
content("O pedido · Board PMI", "O que pedimos ao board", b_ask_board,
    "Variacao 1 de 3 (trocavel por audiencia). Para o board: endorsement estrategico e reconhecimento do "
    "Nucleo como leitura pratica de PMI:Next / M.O.R.E. Ancorar no fato de que a integracao de silos ja e "
    "institucional; o Nucleo e a expressao humana e de pesquisa dela. Manter as outras duas variacoes "
    "ocultas conforme a plateia.")

def b_ask_chapter(s, top):
    add_box(s, 0.62, top, 12.1, 1.7, "Para os presidentes de capítulo", [
        "Adesão do capítulo à horizontal Núcleo IA & GP (federação leve, sem perder identidade local).",
        "Indicação de protagonistas, pesquisadores e curadores do capítulo para as verticais."],
        hcolor=PURPLE)
    add_box(s, 0.62, top+1.9, 12.1, 1.6, "O que o capítulo ganha", [
        "Acesso à produção de pesquisa e à escada Champion → CPMAI para os seus membros.",
        "Validação do discurso AI-forward localmente, alinhado à estratégia global do PMI."],
        hcolor=TEAL)
content("O pedido · Presidentes de capítulo", "O que pedimos aos capítulos", b_ask_chapter,
    "Variacao 2 de 3. Para presidentes de capitulo: adesao do capitulo, indicacao de "
    "protagonistas/pesquisadores e validacao do discurso. Enfatizar que e federacao leve, ganho de "
    "pesquisa e escada de credencial para os membros, sem perder identidade local.")

def b_ask_partner(s, top):
    add_box(s, 0.62, top, 12.1, 1.7, "Para os parceiros de vertical (GPM, Construction Ambassadors, PMOGA)", [
        "Co-curadoria da vertical: o conteúdo fala a credencial, a dor e a linguagem da comunidade.",
        "Acesso à comunidade da credencial para formar a coorte fundadora."],
        hcolor=PURPLE)
    add_box(s, 0.62, top+1.9, 12.1, 1.6, "A regra que protege o parceiro", [
        "Produção é da tribo; a vertical é canal. Nenhuma vertical é dona do conhecimento: é doca.",
        "O pitch de cada vertical se refina com o parceiro antes de ir a público."],
        hcolor=TEAL)
content("O pedido · Parceiros de vertical", "O que pedimos aos parceiros", b_ask_partner,
    "Variacao 3 de 3. Para parceiros de vertical: co-curadoria mais acesso a comunidade da credencial "
    "para a coorte fundadora. Reforcar o principio anti-silo (producao da tribo, distribuicao da "
    "vertical) como a garantia de que o parceiro nao perde dominio sobre o proprio conhecimento.")

# ---- 13 PRÓXIMOS PASSOS (landing §4 + kit ordem de ativação) ----
def b_next(s, top):
    add_table(s, [
        ["Ordem de ativação", "Por que primeiro", "Status"],
        ["1 · Construção", "Líder aceito (Henrique Diniz) + 2 Global Ambassadors (BR/EUA): a mais pronta para ativar", "em formação · líder definido"],
        ["2 · PMO", "PMI-PMOCP recém-lançada + PMOGA absorvida: máximo alinhamento institucional", "em formação · Ciclo 4"],
        ["3 · ESG", "CSPP recém-lançada (janela de mídia) + #1 preditor de sucesso", "em formação · Ciclo 4"],
        ["4 · Ágil", "Agile Alliance no PMI + refresh do PMP (jul)", "declarada"],
        ["5 · Negócio", "definir parceiro primeiro", "declarada"],
    ], top, widths=[1.8, 4.8, 1.9], fs=12)
    add_box(s, 0.62, top+2.95, 12.1, 1.1, "Seja protagonista", [
        "Cada vertical entra com status explícito (em formação, não vaporware). O CTA recruta fundadores, "
        "não consumidores: um programa de liderança, alinhado a M.O.R.E. e PMI:Next."],
        hcolor=ORANGE, hfs=15, fs=12.5)
content("Próximos passos", "Ciclo 4: verticais-piloto e chamada de protagonistas", b_next,
    "Fechamento. A ordem de ativacao agora abre pela PRONTIDAO: Construcao e #1 porque e a unica com "
    "lider aceito (Henrique Diniz) e dois Global Ambassadors (BR/EUA) ja a bordo, com trilha Giga "
    "Projetos e IA. Depois PMO (alinhamento institucional), ESG (janela de midia CSPP, perecivel), Agil "
    "(comunidade se reposicionando), Negocio (definir parceiro). "
    "Cada vertical aparece com status explicito (em formacao), sem fingir atividade. A chamada e 'Seja "
    "protagonista', nao 'seja membro': recruta coorte fundadora, com vinculo a M.O.R.E. e PMI:Next que "
    "justifica chamar de lideranca. Sem hardcode: a pagina le o status da iniciativa e renderiza o CTA.")

# delete template originals
for sldId in SRC: prs.slides._sldIdLst.remove(sldId)

# ===================== guards (fail the build, don't ship it) =====================
def _walk_text(shapes, where, fail):
    for sh in shapes:
        if sh.shape_type == 6:
            _walk_text(sh.shapes, where, fail); continue
        if getattr(sh, "has_table", False):
            for row in sh.table.rows:
                for cell in row.cells:
                    if "—" in cell.text_frame.text: fail(f"em-dash in {where}: table of {sh.name!r}")
            continue
        if sh.has_text_frame and "—" in sh.text_frame.text:
            fail(f"em-dash in {where}: shape {sh.name!r}")

def _assert_guards(p):
    def fail(msg): raise AssertionError(msg)
    sw, sh_ = p.slide_width, p.slide_height
    for i, sl in enumerate(p.slides, 1):
        _walk_text(sl.shapes, f"slide {i}", fail)
        if sl.has_notes_slide and "—" in (sl.notes_slide.notes_text_frame.text or ""):
            fail(f"em-dash in notes of slide {i}")
        for sp in sl.shapes:
            if sp.left is None: continue
            if sp.left + sp.width > sw + 9525 or sp.top + sp.height > sh_ + 9525:
                fail(f"slide {i}: shape {sp.name!r} overflows the canvas")
            if sp.name in BOILER: fail(f"slide {i}: boilerplate {sp.name!r} survived the wipe")
        dividers = [sp for sp in sl.shapes
                    if not sp.has_text_frame and sp.height < 18288 and sp.width > Inches(8)]
        for d in dividers:
            for (slide, l, t, w, h) in INJECTED:
                if slide is not sl: continue
                if t < d.top < t + h and not (l + w < d.left or d.left + d.width < l):
                    fail(f"slide {i}: divider {d.name!r} crosses injected content")
_assert_guards(prs)

prs.save(str(OUT))
print(f"saved {OUT} with {len(prs.slides._sldIdLst)} slides")

# ===================== render: PDF + preview PNGs from the FINAL artifact =====================
subprocess.run(["soffice", "--headless", "--convert-to", "pdf", "--outdir", str(BASE), str(OUT)],
               check=True, capture_output=True)
pdf = OUT.with_suffix(".pdf")
PREVIEW.mkdir(exist_ok=True)
for old in PREVIEW.glob("slide-*.png"): old.unlink()
shutil.copy2(pdf, PREVIEW / pdf.name)
subprocess.run(["pdftoppm", "-png", "-r", "100", str(pdf), str(PREVIEW / "slide")],
               check=True, capture_output=True)
n = len(list(PREVIEW.glob("slide-*.png")))
print(f"rendered {pdf.name} + {n} preview PNGs (from the final PDF)")
assert n == len(prs.slides._sldIdLst), "preview count != slide count"
